"""
Oren Learning Pods -- Markdown-to-PPTX Generator

Converts a pod's README.md into a dark-themed 16:9 PowerPoint presentation.

Usage
-----
    python shared/pptx_generator.py pods/03-data-whisperer/README.md

Output is saved as README.pptx in the same directory as the input file.
"""

from __future__ import annotations

import os
import re
import sys
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# Resolve the shared/ directory so imports work regardless of cwd
_SHARED_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(_SHARED_DIR))

from theme import (
    BG_DARK,
    BG_CARD,
    TEXT_LIGHT,
    TEXT_WHITE,
    OREN_BLUE,
    OREN_GREEN,
    HIGHLIGHT,
    SUCCESS,
    FONT_HEADING,
    FONT_BODY,
    make_rgb_color,
    make_slide_bg,
    hex_to_rgb,
)

# ---------------------------------------------------------------------------
# Slide dimensions -- 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

# ---------------------------------------------------------------------------
# Layout constants (in inches, converted to Emu at point of use)
# ---------------------------------------------------------------------------
MARGIN_LEFT = 0.8
MARGIN_TOP = 0.4
CONTENT_WIDTH = SLIDE_WIDTH_IN - 2 * MARGIN_LEFT
TITLE_HEIGHT = 1.0
BODY_TOP = 1.8
BODY_HEIGHT = SLIDE_HEIGHT_IN - BODY_TOP - 0.5

# Maximum bullet points per content slide before splitting
MAX_BULLETS = 6

# Font for code blocks
FONT_CODE = "Consolas"


# ===================================================================
# Presentation-level state
# ===================================================================

class PptxBuilder:
    """Wraps a python-pptx Presentation with dark-theme slide helpers."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_WIDTH_IN)
        self.prs.slide_height = Inches(SLIDE_HEIGHT_IN)
        self._blank_layout = self._find_blank_layout()

    # ---------------------------------------------------------------
    # Internal helpers
    # ---------------------------------------------------------------

    def _find_blank_layout(self):
        for layout in self.prs.slide_layouts:
            if layout.name == "Blank":
                return layout
        return self.prs.slide_layouts[6]

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self._blank_layout)
        make_slide_bg(slide, BG_DARK)
        return slide

    @staticmethod
    def _add_textbox(slide, left_in, top_in, width_in, height_in):
        return slide.shapes.add_textbox(
            Inches(left_in), Inches(top_in),
            Inches(width_in), Inches(height_in),
        )

    @staticmethod
    def _set_run(run, text, font_name, size_pt, color_hex, bold=False, italic=False):
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.color.rgb = make_rgb_color(color_hex)
        run.font.bold = bold
        run.font.italic = italic

    def _set_paragraph(self, p, text, font_name=FONT_BODY, size_pt=18,
                       color=TEXT_LIGHT, bold=False, italic=False,
                       alignment=PP_ALIGN.LEFT, space_after_pt=0):
        """Write *text* into paragraph *p* with full formatting."""
        p.clear()
        run = p.add_run()
        self._set_run(run, text, font_name, size_pt, color, bold, italic)
        p.alignment = alignment
        if space_after_pt:
            p.space_after = Pt(space_after_pt)

    # ---------------------------------------------------------------
    # Accent bar (thin colored rectangle below a heading)
    # ---------------------------------------------------------------

    def _add_accent_bar(self, slide, top_in, color_hex=HIGHLIGHT, width_in=3.0):
        """Draw a thin accent bar centered horizontally."""
        bar_height = 0.06
        left = (SLIDE_WIDTH_IN - width_in) / 2
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(left), Inches(top_in),
            Inches(width_in), Inches(bar_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = make_rgb_color(color_hex)
        shape.line.fill.background()

    # ---------------------------------------------------------------
    # Slide builders
    # ---------------------------------------------------------------

    def add_title_slide(self, title: str, subtitle: str | None = None):
        """H1 -> large centered title on dark background."""
        slide = self._new_slide()

        # Title
        txb = self._add_textbox(slide, 1.0, 2.2, SLIDE_WIDTH_IN - 2.0, 2.0)
        tf = txb.text_frame
        tf.word_wrap = True
        self._set_paragraph(
            tf.paragraphs[0], title,
            font_name=FONT_HEADING, size_pt=40, color=TEXT_WHITE,
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        # Vertical center
        txb.text_frame.auto_size = None

        # Accent bar
        self._add_accent_bar(slide, 4.4, HIGHLIGHT, 4.0)

        # Subtitle (from blockquote right after H1)
        if subtitle:
            stb = self._add_textbox(slide, 1.5, 4.7, SLIDE_WIDTH_IN - 3.0, 1.0)
            stf = stb.text_frame
            stf.word_wrap = True
            self._set_paragraph(
                stf.paragraphs[0], subtitle,
                font_name=FONT_BODY, size_pt=22, color=HIGHLIGHT,
                italic=True, alignment=PP_ALIGN.CENTER,
            )

        return slide

    def add_section_slide(self, heading: str):
        """H2 -> section divider with centered text and accent bar."""
        slide = self._new_slide()

        txb = self._add_textbox(slide, 1.0, 2.8, SLIDE_WIDTH_IN - 2.0, 1.5)
        tf = txb.text_frame
        tf.word_wrap = True
        self._set_paragraph(
            tf.paragraphs[0], heading,
            font_name=FONT_HEADING, size_pt=32, color=TEXT_WHITE,
            bold=True, alignment=PP_ALIGN.CENTER,
        )
        self._add_accent_bar(slide, 4.5, HIGHLIGHT, 3.0)

        return slide

    def add_content_slide(self, title: str, bullets: list[str]):
        """H3 heading + bullet list.  Splits across slides if > MAX_BULLETS."""
        chunks = [bullets[i:i + MAX_BULLETS] for i in range(0, max(len(bullets), 1), MAX_BULLETS)]

        for idx, chunk in enumerate(chunks):
            slide = self._new_slide()

            # Title
            display_title = title if idx == 0 else f"{title} (cont.)"
            ttb = self._add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, TITLE_HEIGHT)
            tf = ttb.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            self._set_paragraph(
                tf.paragraphs[0], display_title,
                font_name=FONT_HEADING, size_pt=28, color=TEXT_WHITE, bold=True,
            )

            # Accent bar below title
            self._add_accent_bar(slide, MARGIN_TOP + TITLE_HEIGHT + 0.05, HIGHLIGHT, 2.0)

            # Bullets
            if not chunk:
                continue
            btb = self._add_textbox(slide, MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, BODY_HEIGHT)
            btf = btb.text_frame
            btf.word_wrap = True
            btf.auto_size = None

            for i, bullet_text in enumerate(chunk):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                self._add_formatted_bullet(p, bullet_text)

    def _add_formatted_bullet(self, p, text: str, indent_level: int = 0):
        """
        Render a single bullet paragraph.  Handles **bold** fragments
        by splitting into multiple runs with different formatting.
        """
        p.clear()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

        # Indent for nested bullets
        prefix = "\u2022  " if indent_level == 0 else "    \u2013  "
        size = 18 if indent_level == 0 else 16

        # Split on **bold** markers
        parts = re.split(r"(\*\*.+?\*\*)", text)
        first = True
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            is_bold_part = part.startswith("**") and part.endswith("**")
            clean = part.strip("*") if is_bold_part else part
            if first:
                clean = prefix + clean
                first = False
            self._set_run(
                run, clean,
                font_name=FONT_BODY, size_pt=size,
                color_hex=HIGHLIGHT if is_bold_part else TEXT_LIGHT,
                bold=is_bold_part,
            )

    def add_numbered_slide(self, title: str, items: list[str]):
        """Numbered list items under a heading."""
        chunks = [items[i:i + MAX_BULLETS] for i in range(0, max(len(items), 1), MAX_BULLETS)]

        for ci, chunk in enumerate(chunks):
            slide = self._new_slide()
            display_title = title if ci == 0 else f"{title} (cont.)"
            ttb = self._add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, TITLE_HEIGHT)
            tf = ttb.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            self._set_paragraph(
                tf.paragraphs[0], display_title,
                font_name=FONT_HEADING, size_pt=28, color=TEXT_WHITE, bold=True,
            )
            self._add_accent_bar(slide, MARGIN_TOP + TITLE_HEIGHT + 0.05, HIGHLIGHT, 2.0)

            if not chunk:
                continue
            btb = self._add_textbox(slide, MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, BODY_HEIGHT)
            btf = btb.text_frame
            btf.word_wrap = True
            btf.auto_size = None

            base_num = ci * MAX_BULLETS
            for i, item_text in enumerate(chunk):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.clear()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)
                num_run = p.add_run()
                self._set_run(num_run, f"{base_num + i + 1}.  ",
                              FONT_HEADING, 18, HIGHLIGHT, bold=True)
                txt_run = p.add_run()
                self._set_run(txt_run, item_text, FONT_BODY, 18, TEXT_LIGHT)

    def add_callout_to_last_slide(self, text: str):
        """Add a blockquote/callout textbox to the most recently created slide."""
        if not self.prs.slides:
            return
        slide = self.prs.slides[-1]
        ctb = self._add_textbox(slide, MARGIN_LEFT + 0.3, BODY_TOP + BODY_HEIGHT - 0.8,
                                CONTENT_WIDTH - 0.6, 0.7)
        ctf = ctb.text_frame
        ctf.word_wrap = True
        self._set_paragraph(
            ctf.paragraphs[0], f"\u201C{text}\u201D",
            font_name=FONT_BODY, size_pt=16, color=HIGHLIGHT,
            italic=True, alignment=PP_ALIGN.LEFT,
        )

    def add_code_slide(self, title: str, code: str, language: str = ""):
        """Code block -> monospace text on a dark card background."""
        slide = self._new_slide()

        # Title
        ttb = self._add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, TITLE_HEIGHT)
        tf = ttb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        lang_label = f"  ({language})" if language else ""
        self._set_paragraph(
            tf.paragraphs[0], f"{title}{lang_label}",
            font_name=FONT_HEADING, size_pt=28, color=TEXT_WHITE, bold=True,
        )

        # Code card background
        code_top = BODY_TOP - 0.2
        code_height = BODY_HEIGHT + 0.2
        from pptx.enum.shapes import MSO_SHAPE
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(MARGIN_LEFT - 0.1), Inches(code_top),
            Inches(CONTENT_WIDTH + 0.2), Inches(code_height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = make_rgb_color(BG_CARD)
        card.line.fill.background()
        # Subtle corner rounding
        card.adjustments[0] = 0.02

        # Code text -- truncate long code blocks to fit the slide
        lines = code.rstrip("\n").split("\n")
        max_lines = 18
        if len(lines) > max_lines:
            lines = lines[:max_lines] + ["    ..."]
        code_text = "\n".join(lines)

        ctb = self._add_textbox(slide, MARGIN_LEFT + 0.2, code_top + 0.2,
                                CONTENT_WIDTH - 0.2, code_height - 0.4)
        ctf = ctb.text_frame
        ctf.word_wrap = True
        ctf.auto_size = None

        for i, line in enumerate(code_text.split("\n")):
            p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            p.clear()
            run = p.add_run()
            self._set_run(run, line, FONT_CODE, 14, TEXT_LIGHT)
            p.space_after = Pt(2)

    def add_table_slide(self, title: str, headers: list[str], rows: list[list[str]]):
        """Render a markdown table as a styled pptx table."""
        slide = self._new_slide()

        # Title
        ttb = self._add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, TITLE_HEIGHT)
        tf = ttb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        self._set_paragraph(
            tf.paragraphs[0], title,
            font_name=FONT_HEADING, size_pt=28, color=TEXT_WHITE, bold=True,
        )

        # Truncate very large tables
        max_rows = 12
        display_rows = rows[:max_rows]

        n_cols = len(headers)
        n_rows = len(display_rows) + 1  # +1 for header

        table_top = BODY_TOP
        table_height = min(BODY_HEIGHT, 0.45 * n_rows)
        col_width = CONTENT_WIDTH / n_cols

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(MARGIN_LEFT), Inches(table_top),
            Inches(CONTENT_WIDTH), Inches(table_height),
        )
        table = table_shape.table

        # Style header row
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = hdr
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_HEADING
                p.font.size = Pt(14)
                p.font.color.rgb = make_rgb_color(TEXT_WHITE)
                p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = make_rgb_color(OREN_BLUE)

        # Style body rows
        for ri, row_data in enumerate(display_rows):
            for ci in range(n_cols):
                cell_text = row_data[ci] if ci < len(row_data) else ""
                cell = table.cell(ri + 1, ci)
                cell.text = cell_text
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_BODY
                    p.font.size = Pt(12)
                    p.font.color.rgb = make_rgb_color(TEXT_LIGHT)
                bg = BG_CARD if ri % 2 == 0 else BG_DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = make_rgb_color(bg)

    def add_metadata_slide(self, metadata_text: str):
        """Render bold metadata lines (Week / Date / Duration) below the title."""
        if not self.prs.slides:
            return
        slide = self.prs.slides[-1]
        mtb = self._add_textbox(slide, 1.5, 5.0, SLIDE_WIDTH_IN - 3.0, 1.0)
        mtf = mtb.text_frame
        mtf.word_wrap = True
        # Strip leading/trailing bold markers for cleaner display
        clean = metadata_text.replace("**", "")
        self._set_paragraph(
            mtf.paragraphs[0], clean,
            font_name=FONT_BODY, size_pt=16, color=TEXT_LIGHT,
            alignment=PP_ALIGN.CENTER,
        )

    def save(self, path: str):
        self.prs.save(path)


# ===================================================================
# Markdown parser
# ===================================================================

def parse_markdown(md_text: str) -> list[dict]:
    """
    Parse markdown into a list of slide-instruction dicts.

    Each dict has a "type" key and additional payload keys depending on type:
        title     : {text, subtitle}
        section   : {text}
        content   : {title, bullets}
        numbered  : {title, items}
        callout   : {text}
        code      : {title, code, language}
        table     : {title, headers, rows}
        metadata  : {text}
    """
    lines = md_text.split("\n")
    slides: list[dict] = []
    i = 0

    current_h3_title: str | None = None
    current_bullets: list[str] = []
    current_numbered: list[str] = []
    in_code_block = False
    code_lang = ""
    code_lines: list[str] = []
    last_heading_for_code = ""

    def flush_content():
        nonlocal current_h3_title, current_bullets, current_numbered
        if current_h3_title and current_bullets:
            slides.append({"type": "content", "title": current_h3_title, "bullets": current_bullets})
        elif current_h3_title and current_numbered:
            slides.append({"type": "numbered", "title": current_h3_title, "items": current_numbered})
        elif current_h3_title:
            # H3 with no bullets -- create a section-style slide
            slides.append({"type": "section", "text": current_h3_title})
        current_h3_title = None
        current_bullets = []
        current_numbered = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # ----- Code fence -----
        if stripped.startswith("```"):
            if not in_code_block:
                in_code_block = True
                code_lang = stripped[3:].strip()
                code_lines = []
                i += 1
                continue
            else:
                # End of code block
                in_code_block = False
                code_title = last_heading_for_code or "Code"
                slides.append({
                    "type": "code",
                    "title": code_title,
                    "code": "\n".join(code_lines),
                    "language": code_lang,
                })
                code_lang = ""
                code_lines = []
                i += 1
                continue

        if in_code_block:
            code_lines.append(line)
            i += 1
            continue

        # ----- H1 -----
        if stripped.startswith("# ") and not stripped.startswith("## "):
            flush_content()
            title_text = stripped[2:].strip()
            # Look ahead for blockquote subtitle
            subtitle = None
            if i + 1 < len(lines) and lines[i + 1].strip() == "":
                if i + 2 < len(lines) and lines[i + 2].strip().startswith(">"):
                    subtitle = lines[i + 2].strip().lstrip("> ").strip()
            slides.append({"type": "title", "text": title_text, "subtitle": subtitle})
            last_heading_for_code = title_text
            i += 1
            continue

        # ----- H2 -----
        if stripped.startswith("## ") and not stripped.startswith("### "):
            flush_content()
            h2_text = stripped[3:].strip()
            slides.append({"type": "section", "text": h2_text})
            last_heading_for_code = h2_text
            i += 1
            continue

        # ----- H3 -----
        if stripped.startswith("### "):
            flush_content()
            current_h3_title = stripped[4:].strip()
            last_heading_for_code = current_h3_title
            i += 1
            continue

        # ----- Blockquote (not immediately after H1) -----
        if stripped.startswith(">"):
            quote_text = stripped.lstrip("> ").strip()
            # If this is the subtitle line right after H1, skip (already captured)
            if slides and slides[-1].get("type") == "title" and slides[-1].get("subtitle") == quote_text:
                i += 1
                continue
            slides.append({"type": "callout", "text": quote_text})
            i += 1
            continue

        # ----- Bold metadata line (e.g. **Week 3** | **Date:** ...) -----
        if stripped.startswith("**") and "|" in stripped and not stripped.startswith("| "):
            slides.append({"type": "metadata", "text": stripped})
            i += 1
            continue

        # ----- Table -----
        if stripped.startswith("|") and "|" in stripped[1:]:
            # Collect full table
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            if len(table_lines) >= 2:
                headers = [c.strip() for c in table_lines[0].split("|") if c.strip()]
                # Skip separator line (line 1)
                rows = []
                for tl in table_lines[2:]:
                    row = [c.strip() for c in tl.split("|") if c.strip()]
                    if row:
                        rows.append(row)
                table_title = last_heading_for_code or "Table"
                slides.append({"type": "table", "title": table_title, "headers": headers, "rows": rows})
            continue

        # ----- Unordered bullet -----
        bullet_match = re.match(r"^(\s*)[-*]\s+(.+)", line)
        if bullet_match:
            indent = len(bullet_match.group(1))
            text = bullet_match.group(2).strip()
            if indent >= 2 and current_bullets:
                # Nested bullet -- append to last bullet
                current_bullets[-1] += f"  \u2014 {text}"
            else:
                current_bullets.append(text)
            i += 1
            continue

        # ----- Numbered list -----
        num_match = re.match(r"^\s*\d+[.)]\s+(.+)", stripped)
        if num_match:
            current_numbered.append(num_match.group(1).strip())
            i += 1
            continue

        # Skip blank / unrecognized lines
        i += 1

    flush_content()
    return slides


# ===================================================================
# Orchestrator: instructions -> slides
# ===================================================================

def build_presentation(instructions: list[dict]) -> PptxBuilder:
    """Turn parsed markdown instructions into a PptxBuilder with all slides."""
    builder = PptxBuilder()

    for instr in instructions:
        t = instr["type"]

        if t == "title":
            builder.add_title_slide(instr["text"], instr.get("subtitle"))

        elif t == "section":
            builder.add_section_slide(instr["text"])

        elif t == "content":
            builder.add_content_slide(instr["title"], instr["bullets"])

        elif t == "numbered":
            builder.add_numbered_slide(instr["title"], instr["items"])

        elif t == "callout":
            builder.add_callout_to_last_slide(instr["text"])

        elif t == "code":
            builder.add_code_slide(instr["title"], instr["code"], instr.get("language", ""))

        elif t == "table":
            builder.add_table_slide(instr["title"], instr["headers"], instr["rows"])

        elif t == "metadata":
            builder.add_metadata_slide(instr["text"])

    return builder


# ===================================================================
# CLI entry point
# ===================================================================

def main():
    if len(sys.argv) < 2:
        print("Usage: python pptx_generator.py <path-to-README.md>")
        print("Example: python shared/pptx_generator.py pods/03-data-whisperer/README.md")
        sys.exit(1)

    md_path = Path(sys.argv[1]).resolve()
    if not md_path.is_file():
        print(f"Error: file not found: {md_path}")
        sys.exit(1)

    md_text = md_path.read_text(encoding="utf-8")
    instructions = parse_markdown(md_text)

    if not instructions:
        print(f"Warning: no slide content found in {md_path.name}")
        sys.exit(0)

    builder = build_presentation(instructions)

    out_path = md_path.parent / "README.pptx"
    builder.save(str(out_path))

    n_slides = len(builder.prs.slides)
    print(f"Saved: {out_path}")
    print(f"Total slides: {n_slides}")


if __name__ == "__main__":
    main()
