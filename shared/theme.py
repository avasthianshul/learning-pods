"""
Oren Learning Pods -- Dark Theme Constants & Helpers

Shared color palette, font definitions, and python-pptx helper functions
used by the PPTX generator and any other slide-building scripts.
"""

from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

# ---------------------------------------------------------------------------
# Color palette (dark theme)
# ---------------------------------------------------------------------------
BG_DARK = "#1a1a2e"          # main slide background
BG_CARD = "#16213e"          # card / panel / code-block background
TEXT_LIGHT = "#e0e0e0"       # body text
TEXT_WHITE = "#ffffff"        # headings
OREN_BLUE = "#1E2C56"        # primary accent (brand blue)
OREN_GREEN = "#2B7A4B"       # secondary accent (brand green)
HIGHLIGHT = "#4A90D9"        # light blue highlight
SUCCESS = "#0f9b58"          # success / positive indicator

# Per-category accent colors used for section dividers and badges
CATEGORY_COLORS = {
    "Basics & Claude Code":    "#4A90D9",   # light blue
    "Running Code Locally":    "#2B7A4B",   # green
    "Sharing With the World":  "#E67E22",   # orange
    "Connectivity":            "#9B59B6",   # purple
    "Graduation":              "#E74C3C",   # red
}

# ---------------------------------------------------------------------------
# Font constants
# ---------------------------------------------------------------------------
FONT_HEADING = "Lato"
FONT_BODY = "Lato"

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert a hex color string (e.g. '#1a1a2e') to an (r, g, b) tuple."""
    h = hex_color.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def make_rgb_color(hex_color: str) -> RGBColor:
    """Convert a hex color string to a python-pptx RGBColor."""
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


def make_slide_bg(slide, hex_color: str) -> None:
    """Fill *slide*'s background with a solid color given as a hex string."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = make_rgb_color(hex_color)


def make_text(
    text_frame,
    text: str,
    font_size: int,
    color: str,
    bold: bool = False,
    font_name: str = None,
    alignment=None,
    space_after_pt: int = 0,
):
    """
    Convenience wrapper: set text, font, size, color, and optional bold on
    the *first* paragraph of *text_frame*.  Creates a clean run so prior
    formatting is not inherited.

    Parameters
    ----------
    text_frame : pptx text-frame object
    text       : the string to display
    font_size  : point size (int)
    color      : hex color string, e.g. '#ffffff'
    bold       : whether the text is bold
    font_name  : override font (defaults to FONT_BODY)
    alignment  : pptx PP_ALIGN value, or None for default
    space_after_pt : spacing below the paragraph in points
    """
    from pptx.util import Pt as _Pt

    p = text_frame.paragraphs[0]
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.name = font_name or FONT_BODY
    run.font.size = _Pt(font_size)
    run.font.color.rgb = make_rgb_color(color)
    run.font.bold = bold

    if alignment is not None:
        p.alignment = alignment
    if space_after_pt:
        p.space_after = _Pt(space_after_pt)
